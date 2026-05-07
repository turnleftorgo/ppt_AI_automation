from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def create_base_template(filename="base_template.pptx"):
    # 创建一个新的空演示文稿
    prs = Presentation()

    # ==========================================
    # 第一页：封面 (Title Slide)
    # ==========================================
    title_slide_layout = prs.slide_layouts[0] # 0 通常是封面布局
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    # 获取标题和副标题文本框
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]

    # 填入占位符
    title1.text = "{{Main_Title}}"
    subtitle1.text = "{{Sub_Title}}"
    
    # 稍微设置一下样式，验证后续你的核心引擎能否保留样式
    title1.text_frame.paragraphs[0].font.bold = True
    title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204) # 蓝色

    # ==========================================
    # 第二页：正文内容页 (Title and Content)
    # ==========================================
    bullet_slide_layout = prs.slide_layouts[1] # 1 通常是带项目符号的正文布局
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    
    # 获取标题和正文文本框
    title2 = slide2.shapes.title
    body_shape = slide2.placeholders[1]

    # 填入当页标题占位符
    title2.text = "{{Slide_1_Title}}"
    
    # 填入正文项目符号占位符
    tf = body_shape.text_frame
    tf.text = "{{Bullet_Point_1}}"
    
    p2 = tf.add_paragraph()
    p2.text = "{{Bullet_Point_2}}"
    
    p3 = tf.add_paragraph()
    p3.text = "{{Bullet_Point_3}}"

    # 保存为 pptx 文件
    prs.save(filename)
    print(f"成功！已在当前目录生成测试模板: {filename}")

if __name__ == "__main__":
    create_base_template()