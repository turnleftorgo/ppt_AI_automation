"""
One-time script: add {PLACEHOLDER} markers to Standard_report.pptx.

Run from project root:
  cd /mnt/e/workspace/ppt_AI_automation
  python templates/add_placeholders.py
"""
import os
from pptx import Presentation

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(SCRIPT_DIR, "Standard_report.pptx")

# Title text boxes: shape index -> placeholder text
TITLE_PLACEHOLDERS = {
    0: "{REPORT_TITLE}",      # "文字方塊 3" was "5C Report"
    1: "{REPORT_SUBTITLE}",   # "文字方塊 22" was "5C模板"
}

# (row, col, placeholder_name, new_text)
METADATA_PLACEHOLDERS = [
    (0, 1, "FILE_NAME",   "檔案名：\n{FILE_NAME}"),
    (0, 2, "IPAD_MODEL",  "iPad機型：\n{IPAD_MODEL}"),
    (1, 1, "BUILD",       "Build：\n{BUILD}"),
    (1, 2, "IPAD_TYPE",   "iPad機種：\n{IPAD_TYPE}"),
    (2, 1, "PROCESS",     "製程：\n{PROCESS}"),
    (2, 2, "REPORTER",    "報告人：\n{REPORTER}"),
    (3, 1, "KEYWORDS",    "關鍵字：\n{KEYWORDS}"),
]

CONTENT_PLACEHOLDERS = [
    (4, 1, "ISSUE_DESCRIPTION", "{ISSUE_DESCRIPTION}"),
    (5, 1, "ISSUE_ANALYSIS",    "{ISSUE_ANALYSIS}"),
    (6, 1, "ROOT_CAUSE",        "{ROOT_CAUSE}"),
    (7, 1, "CONTAINMENT",       "{CONTAINMENT}"),
    (8, 1, "CORRECTIVE",        "{CORRECTIVE}"),
]


def set_text_frame_text(tf, text: str):
    """Clear text frame and write new text, preserving first run's font style."""
    font_style = None
    for para in tf.paragraphs:
        for run in para.runs:
            try:
                color_rgb = run.font.color.rgb
            except (AttributeError, TypeError):
                color_rgb = None
            font_style = {
                "size": run.font.size,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "name": run.font.name,
                "color_rgb": color_rgb,
            }
            break
        if font_style:
            break

    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""

    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = text
    else:
        first_para.text = text

    if font_style and first_para.runs:
        run = first_para.runs[0]
        try:
            if font_style.get("size"):
                run.font.size = font_style["size"]
            if font_style.get("bold") is not None:
                run.font.bold = font_style["bold"]
            if font_style.get("name"):
                run.font.name = font_style["name"]
            if font_style.get("color_rgb"):
                run.font.color.rgb = font_style["color_rgb"]
        except Exception:
            pass


def main():
    prs = Presentation(PPTX_PATH)
    slide = prs.slides[0]

    # ── Step 1: Title text boxes ──
    count = 0
    for shape_idx, placeholder_text in TITLE_PLACEHOLDERS.items():
        shape = slide.shapes[shape_idx]
        if shape.has_text_frame:
            set_text_frame_text(shape.text_frame, placeholder_text)
            print(f"  Shape[{shape_idx}] '{shape.name}': set to '{placeholder_text}'")
            count += 1

    # ── Step 2: Table cells ──
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if not table:
        print("ERROR: No table found in slide!")
        return

    for row, col, name, text in METADATA_PLACEHOLDERS + CONTENT_PLACEHOLDERS:
        cell = table.cell(row, col)
        set_text_frame_text(cell.text_frame, text)
        print(f"  [{row},{col}] {name}: set to '{text}'")
        count += 1

    prs.save(PPTX_PATH)
    print(f"\nSaved: {PPTX_PATH}")
    print(f"Total placeholders added: {count}")


if __name__ == "__main__":
    main()
